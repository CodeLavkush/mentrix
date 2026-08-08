from pathlib import Path

from src.services.ocr_service import ocr_service
from src.services.pdf_service import pdf_service
from src.utils.logger import logger


class ExtractorService:
    """
    Unified document text extraction service.
    """

    IMAGE_EXTENSIONS = {
        ".jpg",
        ".jpeg",
        ".png",
        ".bmp",
        ".tif",
        ".tiff",
        ".webp",
        ".gif",
    }

    TEXT_EXTENSIONS = {
        ".txt",
        ".md",
        ".markdown",
    }

    DOCX_EXTENSIONS = {
        ".docx",
    }

    PPTX_EXTENSIONS = {
        ".pptx",
    }

    def extract_text(
        self,
        file_path: Path,
    ) -> str:

        suffix = file_path.suffix.lower()

        logger.info(
            "Extracting text from %s",
            file_path.name,
        )

        if suffix == ".pdf":
            return pdf_service.extract_text(
                file_path
            )

        if suffix in self.IMAGE_EXTENSIONS:
            return ocr_service.extract_text(
                file_path
            )

        if suffix in self.TEXT_EXTENSIONS:
            return self._extract_text_file(
                file_path
            )

        if suffix in self.DOCX_EXTENSIONS:
            return self._extract_docx(
                file_path
            )

        if suffix in self.PPTX_EXTENSIONS:
            return self._extract_pptx(
                file_path
            )

        raise ValueError(
            f"Unsupported file type: {suffix}"
        )

    @staticmethod
    def _extract_text_file(
        file_path: Path,
    ) -> str:

        return file_path.read_text(
            encoding="utf-8",
            errors="ignore",
        )

    @staticmethod
    def _extract_docx(
        file_path: Path,
    ) -> str:

        from docx import Document

        document = Document(file_path)

        return "\n".join(
            paragraph.text.strip()
            for paragraph in document.paragraphs
            if paragraph.text.strip()
        )

    @staticmethod
    def _extract_pptx(
        file_path: Path,
    ) -> str:

        from pptx import Presentation

        presentation = Presentation(
            file_path
        )

        slides = []

        for index, slide in enumerate(
            presentation.slides,
            start=1,
        ):
            texts = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()

                    if text:
                        texts.append(text)

            if texts:
                slides.append(
                    f"Slide {index}\n"
                    + "\n".join(texts)
                )

        return "\n\n".join(slides)


extractor_service = ExtractorService()